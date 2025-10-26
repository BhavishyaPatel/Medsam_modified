from pptx import Presentation
from pptx.util import Inches

# --- Presentation Content ---
slides_content = [
    {
        "layout": "title",
        "title": "The Journey of an Image",
        "subtitle": "How an Attention U-Net Creates a Segmentation Mask\n\nBy: Saipranav"
    },
    {
        "layout": "content",
        "title": "Step 1: The Input Tensor",
        "points": [
            "The process begins with a standard image.",
            "This is represented as a 3D data block, or 'tensor'.",
            "Dimensions: 128 (Height) x 128 (Width) x 3 (RGB Channels)."
        ],
        "visual_note": "INSERT DIAGRAM HERE:\nA 3D cuboid representing the input tensor, labeled '128x128x3'."
    },
    {
        "layout": "content",
        "title": "Step 2: The Encoder - Extracting 'What'",
        "points": [
            "The model deconstructs the image to understand its content.",
            "It does this through a repeating cycle of two key operations:",
            "  - Convolution: To find patterns.",
            "  - Max Pooling: To summarize and downsample.",
            "Result: The data gets smaller in size but deeper in meaning."
        ],
        "visual_note": "INSERT DIAGRAM HERE:\nAn arrow showing the 'Encoder Path' of the U-Net, with tensors getting smaller and deeper."
    },
    {
        "layout": "content",
        "title": "Inside the Encoder: Convolution",
        "points": [
            "Small filters (kernels) slide across the input tensor.",
            "Each filter is a specialized pattern detector (e.g., for edges, curves, textures).",
            "Output: A stack of 'feature maps,' each highlighting where a specific pattern was found.",
            "Example: A 128x128x3 tensor becomes a 128x128x64 tensor."
        ],
        "visual_note": "INSERT ANIMATION/DIAGRAM HERE:\nA visual of a 3x3 kernel sliding over an input grid to produce a feature map."
    },
    {
        "layout": "content",
        "title": "Inside the Encoder: Max Pooling",
        "points": [
            "The goal is to summarize the feature maps and reduce noise.",
            "The model takes 2x2 pixel windows and keeps only the strongest signal (the max value).",
            "Result: The height and width of the tensor are halved.",
            "Example: A 128x128x64 tensor becomes a 64x64x64 tensor."
        ],
        "visual_note": "INSERT DIAGRAM HERE:\nA visual showing a 2x2 grid of numbers being reduced to a single, maximum value."
    },
    {
        "layout": "content",
        "title": "Step 3: The Decoder - Rebuilding 'Where'",
        "points": [
            "The model now reconstructs a precise mask from the abstract features.",
            "This involves working back up the 'U' shape.",
            "The key is the Attention Gate, which intelligently guides the reconstruction."
        ],
        "visual_note": "INSERT DIAGRAM HERE:\nAn arrow showing the 'Decoder Path' of the U-Net, with tensors getting larger and shallower."
    },
    {
        "layout": "content",
        "title": "The Magic Trick: The Attention Gate",
        "points": [
            "The gate takes two inputs:",
            "  1. Low-resolution data with good CONTEXT (from the layer below).",
            "  2. High-resolution data with good DETAIL (from the encoder).",
            "It uses the context to create a 'spotlight' that filters the detailed data, keeping only what's relevant to the target.",
        ],
        "visual_note": "INSERT DIAGRAM HERE:\nA diagram of the Attention Gate with two input arrows ('Context', 'Detail') and one 'Focused Output' arrow."
    },
    {
        "layout": "content",
        "title": "Inside the Decoder: The Reconstruction Block",
        "points": [
            "1. Upsample: A small tensor from below is enlarged (e.g., 16x16 -> 32x32).",
            "2. Attention & Filter: The Attention Gate filters the corresponding high-res skip connection data.",
            "3. Concatenate: The upsampled data and the filtered data are stacked together.",
            "4. Convolve: The combined data is processed to refine the features."
        ],
        "visual_note": "INSERT FLOWCHART HERE:\nA flowchart showing these four steps in sequence."
    },
    {
        "layout": "content",
        "title": "Step 4: The Final Prediction",
        "points": [
            "The final decoder output is a deep tensor (e.g., 128x128x64).",
            "A final 1x1 Convolution collapses this into a single channel (128x128x1).",
            "A Sigmoid function converts this into a 'probability map,' where each pixel's value (0 to 1) is the model's confidence.",
            "Thresholding this map at 0.5 creates the final black and white mask."
        ],
        "visual_note": "INSERT VISUAL HERE:\nA visual showing the final tensor being squeezed into the grayscale probability map, then into the binary mask."
    },
    {
        "layout": "title",
        "title": "Thank You",
        "subtitle": "Questions?"
    }
]

# --- Create the Presentation ---
prs = Presentation()
prs.slide_width = Inches(16)
prs.slide_height = Inches(9)

title_slide_layout = prs.slide_layouts[0]
content_slide_layout = prs.slide_layouts[1] # Title and Content Layout

for slide_data in slides_content:
    if slide_data["layout"] == "title":
        slide = prs.slides.add_slide(title_slide_layout)
        title = slide.shapes.title
        subtitle = slide.placeholders[1]
        title.text = slide_data["title"]
        subtitle.text = slide_data["subtitle"]
    
    elif slide_data["layout"] == "content":
        slide = prs.slides.add_slide(content_slide_layout)
        title = slide.shapes.title
        title.text = slide_data["title"]
        
        content_box = slide.placeholders[1]
        tf = content_box.text_frame
        tf.clear()
        
        if "points" in slide_data:
            for point in slide_data["points"]:
                p = tf.add_paragraph()
                p.text = point
                p.level = 0 if not point.strip().startswith("-") else 1

        if "visual_note" in slide_data:
            if "points" in slide_data:
                 p = tf.add_paragraph(); p.text = "\n"
            
            p = tf.add_paragraph()
            p.text = slide_data["visual_note"]
            p.font.bold = True
            p.font.italic = True

# --- Save the Presentation ---
output_filename = "UNet_Backend_Process.pptx"
prs.save(output_filename)

print(f"Presentation saved successfully as {output_filename}")
